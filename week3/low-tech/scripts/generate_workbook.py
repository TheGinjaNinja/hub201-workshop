"""
Hub201 Week 3 - Low-Tech Workbook Generator
Produces: week3-workbook.pptx
Run: python3 week3/low-tech/scripts/generate_workbook.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# Colour constants
GNV_ORANGE = RGBColor(0xEA, 0x58, 0x0D)
BLACK = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF4, 0xF5)
MID_GREY = RGBColor(0x64, 0x74, 0x8B)
BORDER_GREY = RGBColor(0xE4, 0xE4, 0xE7)

# Slide dimensions: standard widescreen 13.33 x 7.5 inches
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUTPUT_PATH = os.path.join(
    os.path.dirname(__file__), "..", "week3-workbook.pptx"
)


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=None, align=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None,
             line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_orange_eyebrow(slide, text, top=Inches(0.4)):
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.35),
                text, font_size=10, bold=True, color=GNV_ORANGE)


def add_slide_title(slide, title, top=Inches(0.75)):
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.7),
                title, font_size=28, bold=True, color=BLACK)


def add_orange_rule(slide, top=Inches(1.5)):
    rect = add_rect(slide, Inches(0.6), top, Inches(2), Inches(0.04),
                    fill_color=GNV_ORANGE)
    return rect


def add_body(slide, text, top=Inches(1.65), font_size=13, color=None):
    if color is None:
        color = BLACK
    return add_textbox(slide, Inches(0.6), top, Inches(12.1), Inches(5.5),
                       text, font_size=font_size, color=color)


# -----------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Orange accent bar on left
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=GNV_ORANGE)

    # Hub201 eyebrow
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
                "HUB201 WEEK 3", font_size=11, bold=True, color=GNV_ORANGE)

    # Main title
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(11.5), Inches(1.5),
                "Customer Discovery and ICP Workbook",
                font_size=40, bold=True, color=BLACK)

    # Orange rule
    add_rect(slide, Inches(0.5), Inches(3.95), Inches(3), Inches(0.06),
             fill_color=GNV_ORANGE)

    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(4.15), Inches(11), Inches(0.6),
                "Low-tech workbook for teams working without an AI coach",
                font_size=16, color=MID_GREY)

    # Footer
    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                "Fill in by hand. Take your completed workbook into Week 4.",
                font_size=11, color=MID_GREY)


def build_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "HOW TO USE THIS WORKBOOK")
    add_slide_title(slide, "Before you begin")
    add_orange_rule(slide)

    body = (
        "This workbook covers three exercises from Hub201 Week 3:\n\n"
        "  Exercise 1 - Hypothesis mapping: define your PULL hypothesis and "
        "rank your top five assumptions by risk.\n\n"
        "  Exercise 2 - Interview plan: design an interview script that "
        "surfaces the Five Signals, not yes/no answers.\n\n"
        "  Exercise 3 - Contact list: build a credibility-first list of "
        "25 target interviewees.\n\n"
        "Each exercise has three parts:\n"
        "  Context (what and why)  |  Template (fill in)  |  Assessment "
        "(review and next action)\n\n"
        "Work through each template, then complete the assessment before "
        "moving on. Bring the finished workbook to Week 4 cohort review."
    )
    add_body(slide, body, top=Inches(1.7), font_size=13)


def build_ex1_context(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "EXERCISE 1 OF 3")
    add_slide_title(slide, "Hypothesis Mapping: what it is")
    add_orange_rule(slide)

    body = (
        "Most early-stage teams work from product assumptions, not customer "
        "evidence. This exercise forces you to state what you believe before "
        "you go out and test it.\n\n"
        "PULL Hypothesis format (one sentence per line):\n\n"
        "  P - Who has BLOCKED demand? (the customer group)\n"
        "  U - What is the project they are URGENTLY prioritising?\n"
        "  L - What do existing options LACK? (specific gap, not vague)\n"
        "  L - What LEADS them to look for a solution right now "
        "(trigger event)?\n\n"
        "Your top five assumptions are the beliefs your whole business "
        "depends on. Rank them: if any one of these is wrong, what kills "
        "the model fastest?\n\n"
        "Kill criteria: what evidence would prove an assumption false?\n"
        "Validate criteria: what evidence would give you confidence to "
        "proceed?"
    )
    add_body(slide, body, top=Inches(1.7), font_size=13)


def build_ex1_template(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "EXERCISE 1 - TEMPLATE")
    add_slide_title(slide, "Hypothesis Mapping: fill in")
    add_orange_rule(slide)

    # --- Top 5 assumptions table ---
    add_textbox(slide, Inches(0.6), Inches(1.65), Inches(12), Inches(0.35),
                "Your top 5 assumptions (rank 1 = highest risk)",
                font_size=11, bold=True, color=MID_GREY)

    headers = ["#", "Assumption (what you believe)", "Kill criteria",
               "Validate criteria"]
    col_w = [Inches(0.4), Inches(4.4), Inches(3.4), Inches(3.4)]
    col_x = [Inches(0.6), Inches(1.05), Inches(5.5), Inches(8.95)]
    row_h = Inches(0.45)
    hdr_top = Inches(2.05)

    # Header row
    for i, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
        r = add_rect(slide, x, hdr_top, w, row_h,
                     fill_color=LIGHT_GREY, line_color=BORDER_GREY,
                     line_width=Pt(0.75))
        add_textbox(slide, x + Inches(0.05), hdr_top + Inches(0.06),
                    w - Inches(0.1), row_h - Inches(0.1),
                    hdr, font_size=10, bold=True, color=BLACK)

    # 5 data rows
    for row_idx in range(5):
        row_top = hdr_top + row_h * (row_idx + 1)
        for i, (x, w) in enumerate(zip(col_x, col_w)):
            add_rect(slide, x, row_top, w, row_h,
                     fill_color=WHITE, line_color=BORDER_GREY,
                     line_width=Pt(0.75))
            if i == 0:
                add_textbox(slide, x + Inches(0.08), row_top + Inches(0.1),
                            w - Inches(0.15), row_h - Inches(0.1),
                            str(row_idx + 1), font_size=11, bold=True,
                            color=MID_GREY)

    # --- PULL hypothesis block ---
    pull_top = Inches(4.55)
    add_textbox(slide, Inches(0.6), pull_top, Inches(12), Inches(0.35),
                "PULL Hypothesis (complete each line)",
                font_size=11, bold=True, color=MID_GREY)

    pull_labels = [
        "P  Who has blocked demand:",
        "U  What project are they prioritising:",
        "L  What do existing options lack:",
        "L  What trigger event leads them to act:",
    ]
    for i, label in enumerate(pull_labels):
        lbl_top = pull_top + Inches(0.38) + Inches(0.52) * i
        add_textbox(slide, Inches(0.6), lbl_top, Inches(3.2), Inches(0.42),
                    label, font_size=11, bold=False, color=MID_GREY)
        add_rect(slide, Inches(3.85), lbl_top + Inches(0.28),
                 Inches(9.1), Inches(0.02), fill_color=BORDER_GREY)


def build_ex1_assessment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "EXERCISE 1 - ASSESSMENT")
    add_slide_title(slide, "Hypothesis Mapping: self-review")
    add_orange_rule(slide)

    add_textbox(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(0.35),
                "Evidence rating for each assumption (circle one per row)",
                font_size=11, bold=True, color=MID_GREY)

    rating_labels = ["Strong evidence", "Moderate evidence",
                     "Weak evidence", "Pure assumption"]
    rating_colors = [
        RGBColor(0x16, 0xA3, 0x4A),
        RGBColor(0xCA, 0x8A, 0x04),
        RGBColor(0xEA, 0x58, 0x0D),
        RGBColor(0xDC, 0x26, 0x26),
    ]

    hdr_top = Inches(2.1)
    row_h = Inches(0.42)
    col_x_num = Inches(0.6)
    rating_starts = [Inches(1.2), Inches(3.9), Inches(6.6), Inches(9.3)]
    col_w_num = Inches(0.5)
    rating_w = Inches(2.5)

    # Header
    add_textbox(slide, col_x_num, hdr_top, col_w_num, row_h,
                "#", font_size=10, bold=True, color=BLACK)
    for j, (lbl, col) in enumerate(zip(rating_labels, rating_colors)):
        add_textbox(slide, rating_starts[j], hdr_top, rating_w, row_h,
                    lbl, font_size=10, bold=True, color=col)

    for row_idx in range(5):
        row_top = hdr_top + row_h * (row_idx + 1)
        add_rect(slide, col_x_num, row_top, col_w_num, row_h,
                 fill_color=LIGHT_GREY, line_color=BORDER_GREY,
                 line_width=Pt(0.75))
        add_textbox(slide, col_x_num + Inches(0.1),
                    row_top + Inches(0.08), col_w_num - Inches(0.15),
                    row_h - Inches(0.1),
                    str(row_idx + 1), font_size=11, bold=True, color=MID_GREY)
        for j in range(4):
            add_rect(slide, rating_starts[j], row_top, rating_w, row_h,
                     fill_color=WHITE, line_color=BORDER_GREY,
                     line_width=Pt(0.75))

    # Biggest risk + Next action
    bottom = Inches(4.55)
    add_textbox(slide, Inches(0.6), bottom, Inches(5.8), Inches(0.35),
                "Biggest risk (assumption most likely to kill the model):",
                font_size=11, bold=True, color=MID_GREY)
    add_rect(slide, Inches(0.6), bottom + Inches(0.4), Inches(5.6),
             Inches(0.7), fill_color=WHITE, line_color=BORDER_GREY,
             line_width=Pt(0.75))

    add_textbox(slide, Inches(7.1), bottom, Inches(5.8), Inches(0.35),
                "Next action (what you will do before next week):",
                font_size=11, bold=True, color=MID_GREY)
    add_rect(slide, Inches(7.1), bottom + Inches(0.4), Inches(5.6),
             Inches(0.7), fill_color=WHITE, line_color=BORDER_GREY,
             line_width=Pt(0.75))


def build_ex2_context(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "EXERCISE 2 OF 3")
    add_slide_title(slide, "Interview Plan: what it is")
    add_orange_rule(slide)

    body = (
        "A bad interview gets validation. A good interview gets evidence. "
        "The difference is whether you ask about the future or the past.\n\n"
        "The Open Question Rule: use why, what, how. Never ask do, would, "
        "could, should. Future-oriented questions invite polite agreement. "
        "Past-behaviour questions surface real decisions.\n\n"
        "The Five Signals you are listening for:\n\n"
        "  1. Recency and frequency - when did this last happen, how often?\n"
        "  2. Concrete behaviour - what exactly did they do, not plan to do?\n"
        "  3. Time and money spent - how much did they invest in solving it?\n"
        "  4. Problem-area match - is the pain in the domain you are "
        "targeting?\n"
        "  5. Awareness level - do they know they have this problem, or do "
        "they call it something else?\n\n"
        "Your script should cover all five signals. Each section below "
        "gives you blank slots. Tag each question with the signal it "
        "is designed to surface."
    )
    add_body(slide, body, top=Inches(1.7), font_size=12)


def build_ex2_template(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "EXERCISE 2 - TEMPLATE")
    add_slide_title(slide, "Interview Script: fill in")
    add_orange_rule(slide, top=Inches(1.5))

    # 5 columns of question slots
    sections = [
        ("INTRO (3 slots)", Inches(0.35), [
            "Q1:", "Q2:", "Q3:"
        ]),
        ("INDUSTRY (3 slots)", Inches(2.8), [
            "Q4:", "Q5:", "Q6:"
        ]),
        ("BUSINESS (4 slots)", Inches(5.25), [
            "Q7:", "Q8:", "Q9:", "Q10:"
        ]),
        ("PERSONAL (4 slots)", Inches(7.7), [
            "Q11:", "Q12:", "Q13:", "Q14:"
        ]),
        ("OUTRO (3 slots)", Inches(10.15), [
            "Q15:", "Q16:", "Magic wand Q:"
        ]),
    ]

    col_w = Inches(2.3)
    slot_h = Inches(0.95)
    base_top = Inches(1.65)

    for (sec_name, sec_x, slots) in sections:
        # Section header
        add_rect(slide, sec_x, base_top, col_w, Inches(0.3),
                 fill_color=GNV_ORANGE)
        add_textbox(slide, sec_x + Inches(0.05), base_top + Inches(0.02),
                    col_w - Inches(0.1), Inches(0.28),
                    sec_name, font_size=9, bold=True, color=WHITE)

        for row_i, q_label in enumerate(slots):
            q_top = base_top + Inches(0.32) + slot_h * row_i
            add_rect(slide, sec_x, q_top, col_w, slot_h - Inches(0.05),
                     fill_color=WHITE, line_color=BORDER_GREY,
                     line_width=Pt(0.75))
            add_textbox(slide, sec_x + Inches(0.06), q_top + Inches(0.04),
                        col_w - Inches(0.12), Inches(0.22),
                        q_label, font_size=9, bold=True, color=GNV_ORANGE)
            # Signal tag line
            add_textbox(slide, sec_x + Inches(0.06),
                        q_top + slot_h - Inches(0.32),
                        col_w - Inches(0.12), Inches(0.22),
                        "Signal tag: ________________",
                        font_size=8, color=MID_GREY)


def build_ex2_assessment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "EXERCISE 2 - ASSESSMENT")
    add_slide_title(slide, "Interview Script: self-review")
    add_orange_rule(slide)

    add_textbox(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(0.35),
                "Signal coverage check: does your script surface all "
                "five signals? Tick each one covered.",
                font_size=12, color=MID_GREY)

    signals = [
        "Recency and frequency",
        "Concrete behaviour",
        "Time and money spent",
        "Problem-area match",
        "Awareness level",
    ]
    for i, sig in enumerate(signals):
        top = Inches(2.15) + Inches(0.42) * i
        add_rect(slide, Inches(0.6), top, Inches(0.35), Inches(0.32),
                 fill_color=WHITE, line_color=BORDER_GREY, line_width=Pt(1))
        add_textbox(slide, Inches(1.05), top, Inches(5), Inches(0.32),
                    sig, font_size=12, color=BLACK)

    # Gaps
    add_textbox(slide, Inches(0.6), Inches(4.4), Inches(5.5), Inches(0.35),
                "Gaps or weak sections (which section needs more questions?):",
                font_size=11, bold=True, color=MID_GREY)
    add_rect(slide, Inches(0.6), Inches(4.8), Inches(5.5), Inches(0.75),
             fill_color=WHITE, line_color=BORDER_GREY, line_width=Pt(0.75))

    add_textbox(slide, Inches(7.0), Inches(4.4), Inches(5.5), Inches(0.35),
                "Next action (revise script, pilot with a teammate?):",
                font_size=11, bold=True, color=MID_GREY)
    add_rect(slide, Inches(7.0), Inches(4.8), Inches(5.5), Inches(0.75),
             fill_color=WHITE, line_color=BORDER_GREY, line_width=Pt(0.75))


def build_ex3_context(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "EXERCISE 3 OF 3")
    add_slide_title(slide, "Contact List: what it is")
    add_orange_rule(slide)

    body = (
        "Your interview outcomes depend almost entirely on who you talk to. "
        "Talking to the wrong people gives you confident but useless data.\n\n"
        "Credibility-first rule: before anyone agrees to talk to you, they "
        "are deciding whether you are worth their time. You have three levers "
        "to build credibility fast:\n\n"
        "  1. Website - live, professional, explains the problem you solve\n"
        "  2. LinkedIn - updated, shows relevant background, "
        "not just your product\n"
        "  3. Social proof - at least one warm introduction, advisor name, "
        "or reference from someone they know\n\n"
        "ICP discipline: every name on your list should be the same type of "
        "person. If your list has 10 different job titles and 8 different "
        "industries, your ICP is not defined.\n\n"
        "Target: 25 contacts minimum. Aim for 10 booked calls in Week 4."
    )
    add_body(slide, body, top=Inches(1.7), font_size=13)


def build_ex3_template(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "EXERCISE 3 - TEMPLATE")
    add_slide_title(slide, "Contact List: fill in")
    add_orange_rule(slide)

    # ICP one-liner
    add_textbox(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(0.3),
                "ICP definition (one sentence): "
                "_____________________________________________________",
                font_size=12, color=MID_GREY)

    # Credibility checklist
    add_textbox(slide, Inches(0.6), Inches(2.05), Inches(3), Inches(0.28),
                "Credibility checklist:", font_size=11, bold=True,
                color=MID_GREY)
    cred_items = ["Website live and professional",
                  "LinkedIn updated with relevant background",
                  "At least one warm intro or social proof reference"]
    for i, item in enumerate(cred_items):
        top = Inches(2.35) + Inches(0.38) * i
        add_rect(slide, Inches(0.6), top, Inches(0.28), Inches(0.26),
                 fill_color=WHITE, line_color=BORDER_GREY, line_width=Pt(1))
        add_textbox(slide, Inches(0.98), top, Inches(5.5), Inches(0.28),
                    item, font_size=11, color=BLACK)

    # Contact table
    tbl_top = Inches(3.55)
    add_textbox(slide, Inches(0.6), tbl_top - Inches(0.3), Inches(12),
                Inches(0.28),
                "Contact list (20 rows shown - add more on next sheet "
                "or paper):",
                font_size=11, bold=True, color=MID_GREY)

    hdrs = ["Name", "Role", "Company", "Why them", "Channel",
            "Status", "Notes"]
    col_ws = [Inches(1.5), Inches(1.5), Inches(1.6), Inches(2.0),
              Inches(1.2), Inches(1.4), Inches(1.8)]
    col_xs = []
    x_acc = Inches(0.3)
    for w in col_ws:
        col_xs.append(x_acc)
        x_acc += w

    row_h = Inches(0.32)
    # Header
    for hdr, cx, cw in zip(hdrs, col_xs, col_ws):
        add_rect(slide, cx, tbl_top, cw, row_h, fill_color=LIGHT_GREY,
                 line_color=BORDER_GREY, line_width=Pt(0.5))
        add_textbox(slide, cx + Inches(0.04), tbl_top + Inches(0.05),
                    cw - Inches(0.08), row_h - Inches(0.08),
                    hdr, font_size=9, bold=True, color=BLACK)

    for row_i in range(10):
        rt = tbl_top + row_h * (row_i + 1)
        for cx, cw in zip(col_xs, col_ws):
            add_rect(slide, cx, rt, cw, row_h, fill_color=WHITE,
                     line_color=BORDER_GREY, line_width=Pt(0.5))


def build_ex3_assessment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_orange_eyebrow(slide, "EXERCISE 3 - ASSESSMENT")
    add_slide_title(slide, "Contact List: self-review")
    add_orange_rule(slide)

    add_textbox(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(0.35),
                "ICP sharpness check: answer each question below.",
                font_size=12, color=MID_GREY)

    questions = [
        "1. Do all your contacts share the same job title or seniority? "
        "(yes / no)",
        "2. Do they all work in the same industry? (yes / no)",
        "3. Do they all face the same specific problem? (yes / no)",
        "4. If any answer is no, which contacts do not fit and should be "
        "removed or moved to a later list?",
        "5. What is the single strongest warm-intro route you have?",
    ]

    for i, q in enumerate(questions):
        top = Inches(2.15) + Inches(0.72) * i
        add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.3),
                    q, font_size=12, color=BLACK)
        add_rect(slide, Inches(0.6), top + Inches(0.32), Inches(12),
                 Inches(0.3), fill_color=WHITE, line_color=BORDER_GREY,
                 line_width=Pt(0.75))

    add_textbox(slide, Inches(0.6), Inches(5.85), Inches(6), Inches(0.35),
                "Next action:", font_size=11, bold=True, color=MID_GREY)
    add_rect(slide, Inches(0.6), Inches(6.25), Inches(12), Inches(0.6),
             fill_color=WHITE, line_color=BORDER_GREY, line_width=Pt(0.75))


def build_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Orange accent bar
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=GNV_ORANGE)

    add_textbox(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(0.4),
                "HUB201 WEEK 3", font_size=11, bold=True, color=GNV_ORANGE)

    add_textbox(slide, Inches(0.5), Inches(2.1), Inches(11.5), Inches(1.0),
                "What to do with this workbook",
                font_size=34, bold=True, color=BLACK)

    add_rect(slide, Inches(0.5), Inches(3.2), Inches(3), Inches(0.06),
             fill_color=GNV_ORANGE)

    body = (
        "1. Complete all three exercises before Week 4. Incomplete "
        "sections are gaps, not excuses.\n\n"
        "2. Share your completed workbook with your cohort team lead "
        "before the next session.\n\n"
        "3. In Week 4 you will present: your PULL hypothesis, "
        "your script plan, and your first booked interviews.\n\n"
        "4. The only thing that counts as evidence in Week 4 is a "
        "confirmed meeting or a completed conversation.\n\n"
        "Good luck. Your goal is 10 conversations, not 10 slides."
    )
    add_textbox(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(3.5),
                body, font_size=14, color=BLACK)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs)
    build_intro(prs)
    build_ex1_context(prs)
    build_ex1_template(prs)
    build_ex1_assessment(prs)
    build_ex2_context(prs)
    build_ex2_template(prs)
    build_ex2_assessment(prs)
    build_ex3_context(prs)
    build_ex3_template(prs)
    build_ex3_assessment(prs)
    build_closing(prs)

    out = os.path.abspath(OUTPUT_PATH)
    prs.save(out)
    print(f"Saved: {out}  ({os.path.getsize(out):,} bytes)")


if __name__ == "__main__":
    main()
